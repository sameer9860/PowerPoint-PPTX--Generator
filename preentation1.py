from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def create_demo_pptx():
    # Create presentation object
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Python PPTX Demo Presentation"
    subtitle.text = "Generated with python-pptx\n" + datetime.now().strftime("%B %d, %Y")
    
    # Slide 2: Content Slide - Introduction
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Introduction to python-pptx"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Features:"
    
    p = tf.add_paragraph()
    p.text = "• Create PowerPoint presentations programmatically"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Add slides, text, images, charts, and tables"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Format text, shapes, and layouts"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• No Microsoft Office required"
    p.level = 1
    
    # Slide 3: Content with Custom Text Box
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Adding Custom Text Boxes"
    
    # Format title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add content text box
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = "You can position text boxes anywhere on the slide:"
    
    p = content_frame.add_paragraph()
    p.text = "• Use Inches() for precise positioning"
    p.level = 0
    p.font.size = Pt(20)
    
    p = content_frame.add_paragraph()
    p.text = "• Control width and height of boxes"
    p.level = 0
    p.font.size = Pt(20)
    
    p = content_frame.add_paragraph()
    p.text = "• Apply custom formatting to text"
    p.level = 0
    p.font.size = Pt(20)
    
    # Slide 4: Slide with Bullet Points
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Python-pptx Capabilities"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    
    tf.text = "What you can create:"
    
    points = [
        "Professional business presentations",
        "Academic and research presentations",
        "Automated report generation",
        "Data visualization with charts",
        "Image galleries and portfolios",
        "Training materials and tutorials"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 5: Table Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Example Data Table"
    
    # Remove default content placeholder
    content = slide.shapes.placeholders[1]
    sp = content.element
    sp.getparent().remove(sp)
    
    # Define table dimensions and position
    rows = 5
    cols = 3
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    
    # Add table
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    
    # Table headers
    headers = ["Feature", "Difficulty", "Usage Frequency"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
    
    # Table data
    data = [
        ["Create Slides", "Easy", "High"],
        ["Add Images", "Easy", "Medium"],
        ["Format Text", "Medium", "High"],
        ["Create Charts", "Hard", "Medium"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_data
            if i % 2 == 0:  # Alternate row colors
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # Slide 6: Final Slide with Contact Info
    slide_layout = prs.slide_layouts[5]  # Title only layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Thank You!"
    
    # Add text box for additional info
    textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(3))
    tf = textbox.text_frame
    
    p = tf.add_paragraph()
    p.text = "This presentation was generated using:"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "python-pptx Library"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\ngithub.com/sameer9860"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    filename = f"demo_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)
    print(f"✓ Presentation created successfully: {filename}")
    return filename

if __name__ == "__main__":
    try:
        output_file = create_demo_pptx()
        print(f"\nThe presentation contains 6 slides demonstrating:")
        print("1. Title slide")
        print("2. Introduction with bullet points")
        print("3. Custom formatted text boxes")
        print("4. Detailed bullet points")
        print("5. Data table with formatting")
        print("6. Final slide with styled text")
    except Exception as e:
        print(f"Error creating presentation: {e}")